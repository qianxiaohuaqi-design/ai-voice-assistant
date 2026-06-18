from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette (cute pastel)
PINK = RGBColor(0xF9, 0xC5, 0xD1)
PURPLE = RGBColor(0xC8, 0xB6, 0xF6)
CREAM = RGBColor(0xFF, 0xF5, 0xE4)
CORAL = RGBColor(0xFF, 0xB4, 0xA2)
DARK_PURPLE = RGBColor(0x6C, 0x5C, 0x9E)
DARK_TEXT = RGBColor(0x3D, 0x3D, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_PINK = RGBColor(0xFD, 0xE2, 0xE9)
LIGHT_PURPLE = RGBColor(0xDD, 0xD0, 0xFA)
SOFT_YELLOW = RGBColor(0xFF, 0xF0, 0xCD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.1
    return shape

def add_circle(slide, left, top, diameter, fill_color):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_text(slide, text, top=Inches(0.6), font_size=40, color=DARK_PURPLE):
    """Add a centered title."""
    return add_text_box(slide, Inches(1), top, Inches(11.333), Inches(1),
                        text, font_size=font_size, color=color, bold=True,
                        alignment=PP_ALIGN.CENTER, font_name='Segoe UI')

def add_decorative_emojis(slide, emojis, top=Inches(0.3), font_size=24):
    """Add a row of decorative emojis across the top."""
    spacing = Inches(12.333) / (len(emojis) + 1)
    for i, emoji in enumerate(emojis):
        left = spacing * (i + 1) - Inches(0.3)
        add_text_box(slide, left, top, Inches(0.6), Inches(0.6),
                     emoji, font_size=font_size, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 1: Cover — 🦀 Meet Claude Code
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, CREAM)

# Decorative floating emojis
add_decorative_emojis(slide1, ['🦀', '⭐', '☁️', '🐱', '💜', '🌸'], top=Inches(0.4), font_size=28)

# Big center card
card = add_rounded_rect(slide1, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5), WHITE, PURPLE)

# Main title
add_text_box(slide1, Inches(2), Inches(1.8), Inches(9.333), Inches(1.2),
             '🦀 认识 Claude Code', font_size=48, color=DARK_PURPLE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, Inches(2.5), Inches(3.0), Inches(8.333), Inches(0.8),
             '你的 AI 编程小助手～', font_size=28, color=RGBColor(0x9B, 0x8B, 0xC4),
             alignment=PP_ALIGN.CENTER)

# Divider line of emojis
add_text_box(slide1, Inches(2), Inches(3.8), Inches(9.333), Inches(0.6),
             '⚡ ⌨️ 💻 🔍 🐛 ✨', font_size=22, alignment=PP_ALIGN.CENTER)

# Bottom description
add_text_box(slide1, Inches(2.5), Inches(4.5), Inches(8.333), Inches(1.2),
             '命令行里的 AI 伙伴 · 像和聪明的螃蟹一起结对编程 🦀',
             font_size=18, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Footer
add_text_box(slide1, Inches(2), Inches(6.3), Inches(9.333), Inches(0.5),
             'Powered by Anthropic  |  Desktop pet by Clawd on Desk 🐾',
             font_size=13, color=RGBColor(0xAA, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

# Corner decorations - small circles
for x_off in [Inches(0.5), Inches(11.8)]:
    add_circle(slide1, x_off, Inches(5.5), Inches(0.35), PINK)
    add_circle(slide1, x_off, Inches(1.2), Inches(0.25), LIGHT_PURPLE)

# ═══════════════════════════════════════════
# SLIDE 2: What is Claude Code? 💭
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, LIGHT_PINK)

add_decorative_emojis(slide2, ['💭', '🐱', '💬', '🌟', '☁️'], top=Inches(0.3), font_size=24)
add_title_text(slide2, '💭 什么是 Claude Code？', top=Inches(0.5), font_size=38)

# Three speech-bubble cards
card_data = [
    ('🖥️', '终端里的 AI', '在你的命令行里直接对话\n不需要打开网页或 IDE 插件\n随时随地召唤它'),
    ('🧠', '理解你的项目', '自动读取整个代码库\n理解上下文，给出精准建议\n就像最懂你的开发搭档'),
    ('🤝', '结对编程伙伴', '一起写代码、改 Bug\n审查代码、搜索文件\n你说需求，它来干活'),
]

for i, (emoji, title, desc) in enumerate(card_data):
    left = Inches(1 + i * 4)
    card = add_rounded_rect(slide2, left, Inches(1.8), Inches(3.3), Inches(4.5), WHITE, CORAL)
    add_text_box(slide2, left + Inches(0.5), Inches(2.1), Inches(2.3), Inches(0.8),
                 f'{emoji}  {title}', font_size=22, color=DARK_PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide2, left + Inches(0.4), Inches(3.0), Inches(2.5), Inches(3.0),
                 desc, font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide2, Inches(2), Inches(6.5), Inches(9.333), Inches(0.6),
             '🐾 不需要任何配置，打开终端就能开始 ✨',
             font_size=16, color=DARK_PURPLE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 3: Core Superpowers ⌨️
# ═══════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, CREAM)

add_decorative_emojis(slide3, ['⌨️', '⚡', '🔍', '🛠️', '🌐'], top=Inches(0.3), font_size=24)
add_title_text(slide3, '⌨️ 核心超能力', top=Inches(0.5), font_size=38)

powers = [
    ('📝', '代码生成与修改', '用自然语言描述需求\n自动写出高质量代码\n精确定位并修改文件'),
    ('🔍', '智能搜索', '秒搜全项目文件\n正则 + 语义理解\n找到你想要的一切'),
    ('🔄', 'Git 操作', '自动提交、分支管理\n生成规范的 commit\n查看历史与差异'),
    ('🌐', '联网搜索', '搜索最新技术资讯\n抓取网页内容分析\n实时获取答案'),
    ('⚙️', '运行脚本', '执行命令和脚本\n安装依赖包\n调试运行错误'),
    ('📦', '文件管理', '创建、编辑、删除文件\n读写图片和 PDF\n管理项目结构'),
]

for i, (emoji, title, desc) in enumerate(powers):
    col = i % 3
    row = i // 3
    left = Inches(1.2 + col * 4)
    top = Inches(2.0 + row * 2.6)

    card = add_rounded_rect(slide3, left, top, Inches(3.3), Inches(2.2), WHITE, PURPLE)
    add_text_box(slide3, left + Inches(0.3), top + Inches(0.15), Inches(2.7), Inches(0.5),
                 f'{emoji}  {title}', font_size=18, color=DARK_PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, left + Inches(0.4), top + Inches(0.7), Inches(2.5), Inches(1.4),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 4: Why so cute & powerful? 🌟
# ═══════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, LIGHT_PURPLE)

add_decorative_emojis(slide4, ['🌟', '🔌', '🎯', '💝', '🦀'], top=Inches(0.3), font_size=24)
add_title_text(slide4, '🌟 为什么这么可爱又强大？', top=Inches(0.5), font_size=38, color=DARK_PURPLE)

features = [
    ('🔌 插件系统', '安装各种插件扩展能力\nGitHub、Playwright、Context7…\n无缝接入你的工作流'),
    ('🎯 Skills 技能', '代码审查、安全扫描\n深度调研、定时任务\n像给螃蟹穿上不同装备'),
    ('🌐 MCP 协议', '连接外部数据源\nSupabase、API、数据库\n让 AI 触达更多信息'),
    ('🤖 智能代理', '自动拆解复杂任务\n多代理并行工作\n像小螃蟹军团协作'),
]

for i, (title, desc) in enumerate(features):
    left = Inches(1 + i * 3.2)
    card = add_rounded_rect(slide4, left, Inches(1.8), Inches(2.7), Inches(3.5), WHITE, CORAL)
    add_text_box(slide4, left + Inches(0.3), Inches(2.1), Inches(2.1), Inches(0.6),
                 title, font_size=20, color=DARK_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide4, left + Inches(0.3), Inches(2.8), Inches(2.1), Inches(2.2),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Heart decorations
add_text_box(slide4, Inches(3), Inches(5.7), Inches(7.333), Inches(0.5),
             '💜  💜  💜  💜  💜  💜  💜  💜  💜',
             font_size=18, alignment=PP_ALIGN.CENTER)

add_text_box(slide4, Inches(2), Inches(6.3), Inches(9.333), Inches(0.6),
             'Clawd 桌宠陪你一起看着 Claude Code 工作 🦀✨',
             font_size=16, color=DARK_PURPLE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 5: When to call it? 🎮
# ═══════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, CREAM)

add_decorative_emojis(slide5, ['🎮', '🎯', '🐛', '📚', '🤖'], top=Inches(0.3), font_size=24)
add_title_text(slide5, '🎮 什么时候找它帮忙？', top=Inches(0.5), font_size=38)

scenarios = [
    ('🚀', '快速原型开发', '「帮我搭一个 Flask API」\n几分钟从零到跑起来'),
    ('🐛', 'Bug 修复', '「这个报错怎么回事？」\n读日志 → 定位 → 改代码'),
    ('📖', '项目学习', '「这个项目是干嘛的？」\n几秒看懂代码结构'),
    ('🔁', '自动化任务', '「每小时跑一次测试」\n定时循环，不用操心'),
    ('🎨', '写文档/PPT', '「帮我做个演示文稿」\n文字到成品一气呵成'),
    ('🧹', '代码重构', '「帮我优化这段代码」\n更清晰、更高效'),
]

for i, (emoji, title, desc) in enumerate(scenarios):
    col = i % 3
    row = i // 3
    left = Inches(1.2 + col * 4)
    top = Inches(1.9 + row * 2.6)

    card = add_rounded_rect(slide5, left, top, Inches(3.3), Inches(2.2), WHITE, CORAL)
    add_text_box(slide5, left + Inches(0.3), top + Inches(0.15), Inches(2.7), Inches(0.5),
                 f'{emoji}  {title}', font_size=18, color=DARK_PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide5, left + Inches(0.4), top + Inches(0.8), Inches(2.5), Inches(1.3),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 6: Ending — Let's go! ☁️
# ═══════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, LIGHT_PINK)

add_decorative_emojis(slide6, ['☁️', '🦀', '🌸', '🐱', '💜', '✨'], top=Inches(0.4), font_size=28)

# Main card
card = add_rounded_rect(slide6, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.8), WHITE, PURPLE)

add_text_box(slide6, Inches(2), Inches(1.8), Inches(9.333), Inches(1.2),
             '☁️ 一起出发吧～', font_size=48, color=DARK_PURPLE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide6, Inches(2.5), Inches(3.0), Inches(8.333), Inches(0.8),
             '打开终端，输入  claude ，然后说 Hello！', font_size=24,
             color=RGBColor(0x9B, 0x8B, 0xC4), alignment=PP_ALIGN.CENTER)

# Links
add_text_box(slide6, Inches(2.5), Inches(4.0), Inches(8.333), Inches(0.5),
             '🔗  github.com/rullerzhou-afk/clawd-on-desk', font_size=18,
             color=CORAL, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide6, Inches(2.5), Inches(4.5), Inches(8.333), Inches(0.5),
             '🌐  claude.ai/code  |  anthropic.com/claude', font_size=16,
             color=DARK_PURPLE, alignment=PP_ALIGN.CENTER)

# Farewell emoji row
add_text_box(slide6, Inches(2), Inches(5.3), Inches(9.333), Inches(0.8),
             '🦀  👋  💜  🐱  ☁️  ✨  🚀',
             font_size=30, alignment=PP_ALIGN.CENTER)

# Footer
add_text_box(slide6, Inches(2), Inches(6.5), Inches(9.333), Inches(0.5),
             '感谢 Clawd on Desk 带来的可爱灵感 🐾  |  Made with ❤️ + python-pptx',
             font_size=13, color=RGBColor(0xAA, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

# Save
output_path = 'D:/claude/Claude_Code_介绍.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
